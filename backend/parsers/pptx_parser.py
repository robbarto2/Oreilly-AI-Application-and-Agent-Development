from pptx import Presentation
from backend.parsers.base import DocumentParser, ParsedContent
from typing import Dict, Any
import logging

logger = logging.getLogger(__name__)


class PPTXParser(DocumentParser):
    """Parser for PowerPoint (PPTX) files with hierarchy preservation"""

    def supports(self, file_path: str) -> bool:
        return file_path.lower().endswith(('.pptx', '.ppt'))

    async def parse(self, file_path: str) -> ParsedContent:
        """
        Parse PPTX file preserving slide hierarchy and speaker notes.

        Extracts:
        - Slide titles
        - Body text and bullet points
        - Speaker notes (very valuable for context)
        - Tables
        - Slide order
        """
        try:
            prs = Presentation(file_path)
            slides = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_content = self._extract_slide_content(slide, slide_num)
                if slide_content:
                    slides.append(slide_content)

            # Create root document node
            root = ParsedContent(
                title=f"Presentation from {file_path}",
                content=f"PowerPoint presentation with {len(slides)} slides",
                item_type="document",
                children=slides,
                metadata={"total_slides": len(slides)},
                sequence_order=0,
            )

            logger.info(f"Parsed PPTX: {len(slides)} slides extracted from {file_path}")
            return root

        except Exception as e:
            logger.error(f"Error parsing PPTX file {file_path}: {e}")
            raise

    def _extract_slide_content(self, slide, slide_num: int) -> ParsedContent:
        """Extract content from a single slide"""
        title = self._extract_title(slide)
        body_text = self._extract_body_text(slide)
        notes = self._extract_notes(slide)
        tables = self._extract_tables(slide)

        # Combine all text content
        content_parts = []
        if body_text:
            content_parts.append(body_text)
        if tables:
            content_parts.append(f"\n\nTables:\n{tables}")

        content = "\n".join(content_parts) if content_parts else ""

        metadata = {
            "slide_number": slide_num,
            "has_notes": bool(notes),
            "has_tables": bool(tables),
        }

        if notes:
            metadata["speaker_notes"] = notes

        return ParsedContent(
            title=title or f"Slide {slide_num}",
            content=content,
            item_type="slide",
            metadata=metadata,
            sequence_order=slide_num,
        )

    def _extract_title(self, slide) -> str:
        """Extract slide title"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        return ""

    def _extract_body_text(self, slide) -> str:
        """Extract body text from all text boxes except title"""
        body_texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes.title:
                text = shape.text.strip()
                if text:
                    body_texts.append(text)

        return "\n".join(body_texts)

    def _extract_notes(self, slide) -> str:
        """Extract speaker notes (very valuable for understanding context)"""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            return text_frame.text.strip()
        return ""

    def _extract_tables(self, slide) -> str:
        """Extract and format table data"""
        tables_content = []

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                table_text = []

                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    table_text.append(row_text)

                if table_text:
                    tables_content.append("\n".join(table_text))

        return "\n\n".join(tables_content) if tables_content else ""
